from __future__ import annotations

import shutil
from pathlib import Path

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

import document_evidence_mcp.parsers.doc as doc_parser
from document_evidence_mcp.config import Settings
from document_evidence_mcp.parsers.doc import DERIVED_DOCX_NAME, DocConversionError
from document_evidence_mcp.service import DocumentEvidenceService


def test_docx_structural_ingest(
    tmp_path: Path,
    service: DocumentEvidenceService,
) -> None:
    path = tmp_path / "regulation.docx"
    document = Document()
    document.add_heading("Clause 4 — Sicherheitsventil", level=1)
    document.add_paragraph("Der Grenzwert beträgt 8 bar.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Parameter"
    table.cell(0, 1).text = "Wert"
    table.cell(1, 0).text = "Prüfdruck"
    table.cell(1, 1).text = "12 bar"
    document.save(path)

    ingested = service.ingest_document(path)
    heading = service.search_evidence("Sicherheitsventil", document_id=ingested.document_id)
    table_hit = service.search_evidence("Prüfdruck", document_id=ingested.document_id)

    assert heading.hits[0].kind == "heading"
    assert table_hit.hits[0].kind == "table"
    assert table_hit.hits[0].bbox is None


def test_doc_structural_ingest_via_word_conversion(
    tmp_path: Path,
    settings: Settings,
    service: DocumentEvidenceService,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    source = tmp_path / "legacy-regulation.doc"
    source.write_bytes(b"legacy OLE fixture placeholder")
    converted_fixture = tmp_path / "converted-fixture.docx"
    document = Document()
    document.add_heading("Legacy Clause 7", level=1)
    document.add_paragraph("Converted pressure limit is 9 bar.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Parameter"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Legacy test pressure"
    table.cell(1, 1).text = "14 bar"
    document.save(converted_fixture)
    conversion_calls = 0

    def fake_convert(input_path: Path, output_path: Path, *, timeout_seconds: int) -> None:
        nonlocal conversion_calls
        conversion_calls += 1
        assert input_path.suffix == ".doc"
        assert timeout_seconds == 60
        shutil.copyfile(converted_fixture, output_path)

    monkeypatch.setattr(doc_parser, "_convert_doc_to_docx", fake_convert)

    ingested = service.ingest_document(source)
    heading = service.search_evidence("Legacy Clause", document_id=ingested.document_id)
    table_hit = service.search_evidence("Legacy test pressure", document_id=ingested.document_id)

    assert ingested.parser == "word-com+python-docx"
    assert heading.hits[0].kind == "heading"
    assert heading.hits[0].parser == "word-com+python-docx"
    assert heading.hits[0].metadata["source_format"] == "doc"
    assert table_hit.hits[0].kind == "table"
    assert not Path(ingested.artifact_dir, "converted-source.docx").exists()
    assert conversion_calls == 1

    converted_cache = (
        settings.object_root
        / ingested.source_sha256[:2]
        / ingested.source_sha256
        / DERIVED_DOCX_NAME
    )
    assert converted_cache.is_file()
    cached_mtime = converted_cache.stat().st_mtime_ns

    forced = service.ingest_document(source, force=True)

    assert forced.revision == 2
    assert converted_cache.stat().st_mtime_ns == cached_mtime
    assert conversion_calls == 1

    converted_cache.write_bytes(b"corrupted derived cache")
    regenerated = service.ingest_document(source, force=True)

    assert regenerated.revision == 3
    assert conversion_calls == 2
    assert converted_cache.stat().st_size > len(b"corrupted derived cache")


def test_doc_conversion_failure_is_actionable_and_atomic(
    tmp_path: Path,
    settings: Settings,
    service: DocumentEvidenceService,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    source = tmp_path / "broken.doc"
    source.write_bytes(b"not a Word document")

    def fail_conversion(
        _input_path: Path,
        _output_path: Path,
        *,
        timeout_seconds: int,
    ) -> None:
        assert timeout_seconds == 60
        raise DocConversionError("Microsoft Word rejected the legacy document")

    monkeypatch.setattr(doc_parser, "_convert_doc_to_docx", fail_conversion)

    with pytest.raises(DocConversionError, match="Word rejected"):
        service.ingest_document(source)

    assert service.list_documents() == []
    assert list(settings.version_root.iterdir()) == []


def test_doc_conversion_timeout_is_actionable(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    source = tmp_path / "slow.doc"
    destination = tmp_path / "slow.docx"
    source.write_bytes(b"legacy")
    monkeypatch.setattr(doc_parser, "_word_com_diagnostic", lambda: None)

    def time_out(command: list[str], **kwargs: object) -> None:
        assert kwargs["timeout"] == 7
        raise doc_parser.subprocess.TimeoutExpired(command, 7)

    monkeypatch.setattr(doc_parser.subprocess, "run", time_out)

    with pytest.raises(DocConversionError, match="timed out after 7 seconds"):
        doc_parser._convert_doc_to_docx(source, destination, timeout_seconds=7)


def test_doc_conversion_script_does_not_depend_on_optional_os_environment_variable() -> None:
    script = Path(doc_parser.__file__).with_name("convert_doc_to_docx.ps1").read_text(
        encoding="utf-8"
    )

    assert 'if ($env:OS' not in script
    assert "OSVersion.Platform" in script


def test_xlsx_preserves_sheet_rows_and_formula_text(
    tmp_path: Path,
    service: DocumentEvidenceService,
) -> None:
    path = tmp_path / "limits.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Grenzwerte"
    sheet.append(["Parameter", "Wert"])
    sheet.append(["压力上限", 12.5])
    sheet.append(["Doppelt", "=B2*2"])
    workbook.save(path)

    ingested = service.ingest_document(path)
    result = service.search_evidence("压力上限", document_id=ingested.document_id)
    formula = service.search_evidence("B2*2", document_id=ingested.document_id)

    assert result.hits[0].section == "Grenzwerte"
    assert result.hits[0].metadata["row"] == 2
    assert formula.hits


def test_pptx_preserves_slide_and_shape_coordinates(
    tmp_path: Path,
    service: DocumentEvidenceService,
) -> None:
    path = tmp_path / "paper.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    shape.text = "Experimentelle Ergebnisse 2026"
    presentation.save(path)

    ingested = service.ingest_document(path)
    result = service.search_evidence("Experimentelle", document_id=ingested.document_id)

    assert result.hits[0].page == 1
    assert result.hits[0].bbox is not None
    assert result.hits[0].metadata["coordinate_space"] == "presentation_emu"
