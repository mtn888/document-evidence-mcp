from __future__ import annotations

from importlib.metadata import version

from pptx import Presentation

from document_evidence_mcp.chunking import chunk_text, normalize_text
from document_evidence_mcp.models import BBox, EvidenceDraft, ParseResult
from document_evidence_mcp.parsers.base import ParseContext


class PptxParser:
    name = "python-pptx"
    version = version("python-pptx")
    extensions = frozenset({".pptx"})

    def parse(self, context: ParseContext) -> ParseResult:
        presentation = Presentation(context.source_path)
        blocks: list[EvidenceDraft] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            for shape_index, shape in enumerate(slide.shapes, start=1):
                bbox = None
                if all(hasattr(shape, name) for name in ("left", "top", "width", "height")):
                    left = float(shape.left)
                    top = float(shape.top)
                    width = float(shape.width)
                    height = float(shape.height)
                    if width > 0 and height > 0:
                        bbox = BBox(x0=left, y0=top, x1=left + width, y1=top + height)
                if getattr(shape, "has_table", False):
                    rows = [
                        "\t".join(normalize_text(cell.text) for cell in row.cells)
                        for row in shape.table.rows
                    ]
                    text = "\n".join(row for row in rows if row.strip())
                    if text:
                        blocks.append(
                            EvidenceDraft(
                                text=text,
                                kind="table",
                                page=slide_index,
                                bbox=bbox,
                                parser=self.name,
                                metadata={
                                    "coordinate_space": "presentation_emu",
                                    "shape_index": shape_index,
                                },
                            )
                        )
                    continue
                text = normalize_text(getattr(shape, "text", ""))
                if not text:
                    continue
                for chunk in chunk_text(
                    text,
                    context.settings.chunk_chars,
                    context.settings.chunk_overlap,
                ):
                    blocks.append(
                        EvidenceDraft(
                            text=chunk,
                            kind="text",
                            page=slide_index,
                            bbox=bbox,
                            parser=self.name,
                            metadata={
                                "coordinate_space": "presentation_emu",
                                "shape_index": shape_index,
                            },
                        )
                    )
        return ParseResult(
            parser=self.name,
            parser_version=self.version,
            page_count=len(presentation.slides),
            blocks=blocks,
        )
